#!/usr/bin/env python3
"""Generates Tests/PowerpointViewerTests/Fixtures/fixture.pptx.

The fixture deliberately encodes scenarios that broke historically, so the
parser tests fail loudly if any of them regress:

- Slide 1: white background + a table whose text carries black srgbClr values.
  (Canary for the XPath context-escape bug where the background stole the
  table's black — background must stay white.)
- Slide 2: centered bold paragraph (the centered-lyrics regression), plus a
  run with explicit Arial and a whitespace-only run between words (libxml2
  whitespace-drop canary: "kami'y" + " " + "nagpapasalamat").
- Slide 3: bulleted list at two levels, and colored shapes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]

# ── Slide 1: table + white bg (XPath canary) ────────────────────────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid()
s.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
table = s.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(8), Inches(2)).table
headers = ["Awit", "Pages", "Notes"]
for c, text in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = text
    run = cell.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    run.font.underline = True
    run.font.color.rgb = RGBColor(0x7F, 0x5F, 0xBF)
for c, text in enumerate(["8 - Page 5", "10 - Page 8", "ok"]):
    cell = table.cell(1, c)
    cell.text = text
    cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)

# ── Slide 2: centered paragraph + whitespace-only run ───────────────────────
s = prs.slides.add_slide(blank)
tb = s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "kami'y"
r1.font.size = Pt(24); r1.font.bold = True; r1.font.name = "Arial"
r2 = p.add_run(); r2.text = " "          # whitespace-only run (libxml2 canary)
r2.font.size = Pt(24)
r3 = p.add_run(); r3.text = "nagpapasalamat"
r3.font.size = Pt(24); r3.font.bold = True; r3.font.name = "Arial"

# ── Slide 3: bullets + shapes ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(2))
tf = tb.text_frame
tf.text = "First bullet"
p2 = tf.add_paragraph(); p2.text = "Nested bullet"; p2.level = 1
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(3), Inches(1.5))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(3), Inches(2), Inches(2))
oval.fill.solid(); oval.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)

out = os.path.join(os.path.dirname(__file__), "..",
                   "Tests", "PowerpointViewerTests", "Fixtures", "fixture.pptx")
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("wrote", os.path.abspath(out))
